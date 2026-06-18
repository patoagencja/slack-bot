"""PPTX presentation generator — Claude API generates content, python-pptx renders it."""
import io
import json
import logging
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

logger = logging.getLogger(__name__)

# ── Brand palette ──────────────────────────────────────────────────────────────
C_PURPLE = RGBColor(123, 97, 255)
C_PINK   = RGBColor(255, 77, 139)
C_ORANGE = RGBColor(255, 107, 53)
C_TEAL   = RGBColor(0, 191, 165)
C_NAVY   = RGBColor(18, 18, 48)
C_WHITE  = RGBColor(255, 255, 255)
C_LIGHT  = RGBColor(240, 238, 255)
C_GRAY   = RGBColor(100, 100, 120)
C_DARK   = RGBColor(30, 30, 60)

# Slide size: 33.87 x 19.05 cm  (16:9 widescreen)
SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

KPI_COLORS = [C_PURPLE, C_PINK, C_ORANGE, C_TEAL, RGBColor(90, 180, 255), RGBColor(200, 100, 255)]


# ── Helpers ────────────────────────────────────────────────────────────────────

def _bg(slide, color: RGBColor):
    """Fill slide background with a solid color."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _box(slide, x, y, w, h, text="", font_size=18, bold=False,
         color=C_WHITE, align=PP_ALIGN.LEFT, bg=None, alpha=None):
    """Add a text box. Returns the shape."""
    txBox = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    if bg:
        fill = txBox.fill
        fill.solid()
        fill.fore_color.rgb = bg
    return txBox


def _rect(slide, x, y, w, h, color: RGBColor, radius=False):
    """Add a filled rectangle."""
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()  # no border
    return shape


def _add_slide(prs):
    """Add a blank slide."""
    blank_layout = prs.slide_layouts[6]
    return prs.slides.add_slide(blank_layout)


def _multi_para(tf, lines, font_size=14, color=C_WHITE, bold=False, spacing_after=6):
    """Fill a text frame with multiple paragraphs."""
    from pptx.util import Pt as _Pt
    from pptx.oxml.ns import qn
    from lxml import etree
    tf.clear()
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.size = _Pt(font_size)
        run.font.color.rgb = color
        run.font.bold = bold
        # spacing after
        pPr = p._pPr
        if pPr is None:
            pPr = p._p.get_or_add_pPr()
        spcAft = etree.SubElement(pPr, qn('a:spcAft'))
        spcPts = etree.SubElement(spcAft, qn('a:spcPts'))
        spcPts.set('val', str(spacing_after * 100))


# ── Slide builders ─────────────────────────────────────────────────────────────

def _slide_title(prs, title, client_name="", subtitle="", date_range="", agency="Pato Agency"):
    slide = _add_slide(prs)
    _bg(slide, C_NAVY)

    # Accent bar (3 segments, top)
    segment_w = 13.33 / 3
    for i, col in enumerate([C_PURPLE, C_PINK, C_ORANGE]):
        _rect(slide, i * segment_w, 0, segment_w, 0.12, col)

    # Decorative circles (right side)
    for size, alpha_color in [(2.5, RGBColor(123, 97, 255)), (1.5, RGBColor(255, 77, 139))]:
        shape = slide.shapes.add_shape(9, Inches(10.8), Inches(2.5 - size/2), Inches(size), Inches(size))
        shape.fill.solid()
        shape.fill.fore_color.rgb = alpha_color
        shape.line.fill.background()

    # Agency label
    _box(slide, 0.5, 0.3, 5, 0.5, agency.upper(),
         font_size=9, color=C_PURPLE, bold=True)

    # Client badge
    if client_name:
        badge = slide.shapes.add_shape(1, Inches(0.5), Inches(1.0), Inches(2.2), Inches(0.42))
        badge.fill.solid()
        badge.fill.fore_color.rgb = C_PURPLE
        badge.line.fill.background()
        tf = badge.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = client_name.upper()
        run.font.size = Pt(11)
        run.font.bold = True
        run.font.color.rgb = C_WHITE

    # Main title
    _box(slide, 0.5, 1.65, 9.5, 2.2, title,
         font_size=38, bold=True, color=C_WHITE)

    # Subtitle
    if subtitle:
        _box(slide, 0.5, 3.85, 9.5, 0.7, subtitle,
             font_size=18, color=C_LIGHT)

    # Date range (bottom left)
    if date_range:
        _box(slide, 0.5, 6.5, 5, 0.5, f"📅  {date_range}",
             font_size=12, color=C_GRAY)

    # Bottom accent bar
    _rect(slide, 0, 7.38, 13.33, 0.12, C_PURPLE)


def _slide_kpis(prs, title, kpis):
    """KPI card slide. kpis = list of {"label", "value", "sub"?, "icon"?}"""
    slide = _add_slide(prs)
    _bg(slide, C_WHITE)

    # Header bar
    _rect(slide, 0, 0, 13.33, 1.1, C_NAVY)
    _box(slide, 0.5, 0.18, 12, 0.8, title,
         font_size=24, bold=True, color=C_WHITE)

    # Accent strip bottom of header
    _rect(slide, 0, 1.1, 4.44, 0.07, C_PURPLE)
    _rect(slide, 4.44, 1.1, 4.44, 0.07, C_PINK)
    _rect(slide, 8.88, 1.1, 4.45, 0.07, C_ORANGE)

    cols = min(len(kpis), 3)
    rows = (len(kpis) + cols - 1) // cols
    card_w = 12.0 / cols
    card_h = (5.8 / rows) - 0.15
    start_x = 0.67
    start_y = 1.4

    for i, kpi in enumerate(kpis):
        row = i // cols
        col = i % cols
        cx = start_x + col * (card_w + 0.18)
        cy = start_y + row * (card_h + 0.18)

        color = KPI_COLORS[i % len(KPI_COLORS)]

        # Card bg (very light tint)
        card = slide.shapes.add_shape(1, Inches(cx), Inches(cy), Inches(card_w), Inches(card_h))
        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(245, 244, 255)
        card.line.color.rgb = RGBColor(220, 216, 255)
        card.line.width = Pt(1)

        # Top color strip
        _rect(slide, cx, cy, card_w, 0.08, color)

        # Icon + label
        icon = kpi.get("icon", "📊")
        label = kpi.get("label", "")
        _box(slide, cx + 0.15, cy + 0.18, card_w - 0.3, 0.4,
             f"{icon}  {label}", font_size=11, color=C_GRAY, bold=False)

        # Big value
        _box(slide, cx + 0.1, cy + 0.52, card_w - 0.2, card_h - 0.85,
             kpi.get("value", "—"),
             font_size=28, bold=True, color=C_DARK)

        # Sub-label
        if kpi.get("sub"):
            _box(slide, cx + 0.15, cy + card_h - 0.45, card_w - 0.3, 0.4,
                 kpi["sub"], font_size=10, color=C_GRAY)


def _slide_content(prs, title, bullets, header_color=C_NAVY):
    slide = _add_slide(prs)
    _bg(slide, C_WHITE)

    _rect(slide, 0, 0, 13.33, 1.1, header_color)
    _box(slide, 0.5, 0.18, 12, 0.8, title,
         font_size=24, bold=True, color=C_WHITE)
    _rect(slide, 0, 1.1, 0.07, 6.4, C_PURPLE)

    body_box = slide.shapes.add_textbox(Inches(0.4), Inches(1.35), Inches(12.5), Inches(5.8))
    tf = body_box.text_frame
    tf.word_wrap = True

    clean = [b for b in bullets if b.strip()]
    _multi_para(tf, clean, font_size=15, color=C_DARK, spacing_after=8)


def _slide_dark_insights(prs, title, bullets):
    slide = _add_slide(prs)
    _bg(slide, C_NAVY)

    _rect(slide, 0, 0, 0.07, 7.5, C_PURPLE)
    _box(slide, 0.35, 0.3, 12, 0.8, title,
         font_size=28, bold=True, color=C_WHITE)
    _rect(slide, 0.35, 1.15, 3, 0.05, C_PINK)

    body_box = slide.shapes.add_textbox(Inches(0.35), Inches(1.4), Inches(12.5), Inches(5.8))
    tf = body_box.text_frame
    tf.word_wrap = True
    clean = [b for b in bullets if b.strip()]
    _multi_para(tf, clean, font_size=15, color=C_LIGHT, spacing_after=10)


def _slide_table(prs, title, headers, rows):
    from pptx.util import Inches as _I, Pt as _Pt
    slide = _add_slide(prs)
    _bg(slide, C_WHITE)

    _rect(slide, 0, 0, 13.33, 1.1, C_NAVY)
    _box(slide, 0.5, 0.18, 12, 0.8, title,
         font_size=24, bold=True, color=C_WHITE)

    if not rows:
        return

    n_cols = len(headers)
    n_rows = len(rows)
    col_w = 12.0 / n_cols
    row_h = min(0.45, 5.5 / (n_rows + 1))
    table = slide.shapes.add_table(
        n_rows + 1, n_cols,
        _I(0.67), _I(1.3),
        _I(12.0), _I(row_h * (n_rows + 1))
    ).table

    # Header row
    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.fill.solid()
        cell.fill.fore_color.rgb = C_PURPLE
        tf = cell.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        run = tf.paragraphs[0].add_run()
        run.text = str(h)
        run.font.size = _Pt(11)
        run.font.bold = True
        run.font.color.rgb = C_WHITE

    # Data rows
    for ri, row in enumerate(rows):
        bg = RGBColor(245, 244, 255) if ri % 2 == 0 else C_WHITE
        for ci, val in enumerate(row[:n_cols]):
            cell = table.cell(ri + 1, ci)
            cell.fill.solid()
            cell.fill.fore_color.rgb = bg
            tf = cell.text_frame
            tf.paragraphs[0].alignment = PP_ALIGN.CENTER
            run = tf.paragraphs[0].add_run()
            run.text = str(val)
            run.font.size = _Pt(10)
            run.font.color.rgb = C_DARK


# ── Content generation via Claude API ─────────────────────────────────────────

def _generate_slide_structure(title, client_name, subtitle, brief,
                               date_range, meta_ads_data, google_ads_data, extra_slides):
    """Call Claude API to generate structured slide content as JSON."""
    from anthropic import Anthropic
    ai = Anthropic(api_key=os.environ.get("CLAUDE_API_KEY") or os.environ.get("ANTHROPIC_API_KEY"))

    data_ctx = ""
    if meta_ads_data:
        data_ctx += f"\n\nDANE META ADS:\n{json.dumps(meta_ads_data, ensure_ascii=False, indent=2)}"
    if google_ads_data:
        data_ctx += f"\n\nDANE GOOGLE ADS:\n{json.dumps(google_ads_data, ensure_ascii=False, indent=2)}"
    if brief:
        data_ctx += f"\n\nBRIEF:\n{brief}"
    if extra_slides:
        data_ctx += f"\n\nDODATKOWE SLAJDY ZASUGEROWANE:\n{json.dumps(extra_slides, ensure_ascii=False)}"

    prompt = f"""Jesteś ekspertem od prezentacji agencji marketingowej. Masz dane kampanii i musisz wygenerować strukturę prezentacji post-buy (raport wyników dla klienta) jako JSON.

TYTUŁ: {title}
KLIENT: {client_name or ""}
PODTYTUŁ: {subtitle or ""}
DATY: {date_range or ""}
{data_ctx}

Wygeneruj TYLKO JSON (bez markdown, bez ```), struktura:
{{
  "slides": [
    {{
      "type": "title",
      "title": "...",
      "subtitle": "...",
      "client": "...",
      "date_range": "..."
    }},
    {{
      "type": "kpis",
      "title": "Podsumowanie KPI",
      "kpis": [
        {{"label": "Zasięg", "value": "45 232", "icon": "👥", "sub": "unikalnych użytkowników"}},
        {{"label": "Wyświetlenia", "value": "89 450", "icon": "👁️", "sub": "łączne impressions"}},
        {{"label": "Wydatek", "value": "1 234,56 zł", "icon": "💰", "sub": "całkowity budżet"}},
        {{"label": "Kliknięcia", "value": "1 023", "icon": "🖱️", "sub": "link clicks"}},
        {{"label": "CPM", "value": "13,80 zł", "icon": "📊", "sub": "koszt za 1000 wyświetleń"}},
        {{"label": "CPC", "value": "1,21 zł", "icon": "💡", "sub": "koszt za kliknięcie"}}
      ]
    }},
    {{
      "type": "content",
      "title": "Wyniki kampanii",
      "bullets": ["• Kampania X: ...", "• CTR: 2.3%", "..."],
      "dark": false
    }},
    {{
      "type": "table",
      "title": "Wyniki według kampanii",
      "headers": ["Kampania", "Wydatek", "Zasięg", "CTR"],
      "rows": [["Kampania A", "500 zł", "10 000", "2.1%"]]
    }},
    {{
      "type": "content",
      "title": "Wnioski i rekomendacje",
      "bullets": ["✅ ...", "⚠️ ...", "🚀 ..."],
      "dark": true
    }}
  ]
}}

ZASADY:
- Użyj PRAWDZIWYCH liczb z danych (nie szacunków)
- Przelicz zł z USD jeśli spend jest w dolarach (kurs 4.0)
- Polskie opisy, profesjonalne, konkretne
- Slajd wnioski jako ciemny (dark: true)
- Jeśli są dane breakdown (placement, device) → dodaj slajd table z tymi danymi
- Min 4 slajdy, max 8"""

    resp = ai.messages.create(
        model="claude-opus-4-8",
        max_tokens=4000,
        messages=[{"role": "user", "content": prompt}]
    )
    raw = resp.content[0].text.strip()
    raw = re.sub(r"^```[a-z]*\n?", "", raw)
    raw = re.sub(r"\n?```$", "", raw).strip()
    return json.loads(raw)


# ── Main export function ───────────────────────────────────────────────────────

def generate_pptx(title, client_name=None, subtitle=None, brief=None,
                  date_range=None, meta_ads_data=None, google_ads_data=None,
                  extra_slides=None):
    """
    Generate a PPTX presentation and return it as bytes.
    Raises on error.
    """
    try:
        from pptx import Presentation as _Prs  # noqa: verify import
    except ImportError:
        raise RuntimeError("python-pptx nie zainstalowany — uruchom: pip install python-pptx")

    structure = _generate_slide_structure(
        title, client_name, subtitle, brief,
        date_range, meta_ads_data, google_ads_data, extra_slides
    )

    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for slide_def in structure.get("slides", []):
        stype = slide_def.get("type", "content")

        if stype == "title":
            _slide_title(
                prs,
                title=slide_def.get("title", title),
                client_name=slide_def.get("client", client_name or ""),
                subtitle=slide_def.get("subtitle", subtitle or ""),
                date_range=slide_def.get("date_range", date_range or ""),
            )
        elif stype == "kpis":
            _slide_kpis(prs, slide_def.get("title", "KPI"), slide_def.get("kpis", []))
        elif stype == "table":
            _slide_table(
                prs,
                slide_def.get("title", ""),
                slide_def.get("headers", []),
                slide_def.get("rows", []),
            )
        elif stype == "content":
            bullets = slide_def.get("bullets", [])
            dark = slide_def.get("dark", False)
            t = slide_def.get("title", "")
            if dark:
                _slide_dark_insights(prs, t, bullets)
            else:
                _slide_content(prs, t, bullets)

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.read()


def share_pptx(pptx_bytes: bytes, filename: str) -> str | None:
    """
    Upload PPTX bytes to a temporary file host and return a download URL.
    Tries transfer.sh, then 0x0.st as fallback.
    Returns URL string or None on failure.
    """
    import requests as _req
    safe_name = filename.replace(" ", "_")

    # Attempt 1: transfer.sh (14-day links)
    try:
        r = _req.put(
            f"https://transfer.sh/{safe_name}",
            data=pptx_bytes,
            headers={"Max-Days": "14"},
            timeout=30,
        )
        if r.ok and r.text.strip().startswith("http"):
            return r.text.strip()
    except Exception as _e:
        logger.warning(f"transfer.sh failed: {_e}")

    # Attempt 2: 0x0.st (permanent until unused)
    try:
        r = _req.post(
            "https://0x0.st",
            files={"file": (safe_name, pptx_bytes, "application/vnd.openxmlformats-officedocument.presentationml.presentation")},
            timeout=30,
        )
        if r.ok and r.text.strip().startswith("http"):
            return r.text.strip()
    except Exception as _e:
        logger.warning(f"0x0.st failed: {_e}")

    return None
